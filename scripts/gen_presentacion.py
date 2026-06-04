#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Genera la presentacion ejecutiva en .pptx (16:9) con notas de presentador."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INDIGO = RGBColor(0x4F, 0x46, 0xE5)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)
DARK = RGBColor(0x1F, 0x29, 0x37)
GRAY = RGBColor(0x4B, 0x55, 0x63)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xEE, 0xEC, 0xFB)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _notes(slide, txt):
    if txt:
        slide.notes_slide.notes_text_frame.text = txt


def add_band(slide, color, top, height):
    shp = slide.shapes.add_shape(1, 0, top, SW, height)  # 1 = rectangle
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def title_slide(title, subtitle, author):
    s = prs.slides.add_slide(BLANK)
    add_band(s, INDIGO, 0, SH)
    add_band(s, PURPLE, Emu(int(SH * 0.62)), Emu(int(SH * 0.38)))
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = subtitle
    r.font.size = Pt(24); r.font.color.rgb = LIGHT
    tb2 = s.shapes.add_textbox(Inches(0.9), Inches(5.0), Inches(11.5), Inches(1.2))
    p = tb2.text_frame.paragraphs[0]
    r = p.add_run(); r.text = author
    r.font.size = Pt(16); r.font.color.rgb = LIGHT
    return s


def _header(s, title):
    add_band(s, INDIGO, 0, Inches(1.15))
    tb = s.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12.1), Inches(0.85))
    tf = tb.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE


def content_slide(title, bullets, notes=""):
    s = prs.slides.add_slide(BLANK)
    _header(s, title)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(12.0), Inches(5.6))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(10)
        if isinstance(b, tuple):  # (label_bold, rest)
            label, rest = b
            r = p.add_run(); r.text = "•  " + label + " "
            r.font.bold = True; r.font.size = Pt(20); r.font.color.rgb = INDIGO
            r = p.add_run(); r.text = rest
            r.font.size = Pt(20); r.font.color.rgb = DARK
        else:
            r = p.add_run(); r.text = "•  " + b
            r.font.size = Pt(20); r.font.color.rgb = DARK
    _notes(s, notes)
    return s


def mono_slide(title, lines, caption_bullets, notes=""):
    s = prs.slides.add_slide(BLANK)
    _header(s, title)
    box = s.shapes.add_shape(1, Inches(0.7), Inches(1.5), Inches(11.9), Inches(2.5))
    box.fill.solid(); box.fill.fore_color.rgb = DARK; box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
    first = True
    for ln in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        r = p.add_run(); r.text = ln
        r.font.name = "Courier New"; r.font.size = Pt(15); r.font.color.rgb = RGBColor(0x8B, 0xE9, 0xC0)
    tb = s.shapes.add_textbox(Inches(0.7), Inches(4.3), Inches(12.0), Inches(2.8))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for b in caption_bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(8)
        r = p.add_run(); r.text = "•  " + b
        r.font.size = Pt(19); r.font.color.rgb = DARK
    _notes(s, notes)
    return s


def table_slide(title, headers, rows, notes=""):
    s = prs.slides.add_slide(BLANK)
    _header(s, title)
    nrows = len(rows) + 1
    ncols = len(headers)
    gt = s.shapes.add_table(nrows, ncols, Inches(0.9), Inches(1.5),
                            Inches(11.5), Inches(0.5) * nrows).table
    for j, h in enumerate(headers):
        c = gt.cell(0, j); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = INDIGO
        pr = c.text_frame.paragraphs[0]; pr.runs[0].font.bold = True
        pr.runs[0].font.color.rgb = WHITE; pr.runs[0].font.size = Pt(18)
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row):
            c = gt.cell(i, j); c.text = val
            pr = c.text_frame.paragraphs[0]; pr.runs[0].font.size = Pt(16)
            pr.runs[0].font.color.rgb = DARK
            c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
    _notes(s, notes)
    return s


def closing_slide(title, bullets, notes=""):
    s = prs.slides.add_slide(BLANK)
    add_band(s, INDIGO, 0, SH)
    add_band(s, PURPLE, Emu(int(SH * 0.62)), Emu(int(SH * 0.38)))
    tb = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(3.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title
    r.font.size = Pt(44); r.font.bold = True; r.font.color.rgb = WHITE
    for b in bullets:
        p = tf.add_paragraph()
        r = p.add_run(); r.text = b
        r.font.size = Pt(20); r.font.color.rgb = LIGHT
    _notes(s, notes)
    return s


# ---------------- CONTENIDO ----------------
title_slide(
    "Micrositio de Pedidos y Delivery",
    "Proyecto final — Kubernetes / DevOps",
    "Sebastián Contreras · Grupo G-02 · Capa sobre ZUYU (SaaS de punto de venta)",
)

content_slide("El problema y la solución", [
    ("Problema:", "un negocio quiere vender en línea y coordinar entregas sin montar infraestructura."),
    ("Solución:", "app web (catálogo → carrito → pedido → repartidor → entrega) lista para PRODUCCIÓN."),
    "Alta disponibilidad, escalado automático, seguridad de red, CI/CD y observabilidad.",
], "El objetivo no era solo que funcione, sino que esté OPERADA como un sistema de producción real.")

mono_slide("Arquitectura (vista de pájaro)", [
    "Cliente → Ingress (HTTPS) → API (Node/Express)",
    "                              ├── MongoDB (Replica Set ×3)   ← pedidos",
    "                              └── Redis (cola) → Worker → carrier delivery",
], [
    "3 nodos: 1 master + 2 workers (Rocky Linux, kubeadm, CNI Calico).",
    "Todo en el namespace 'micrositio', desplegado por GitOps.",
], "El API atiende al cliente; lo pesado (pedir repartidor, notificar) lo hace un worker asíncrono por una cola, para no bloquear la venta.")

content_slide("MongoDB — Replica Set (alta disponibilidad)", [
    ("Qué es:", "3 instancias MongoDB (StatefulSet) que replican los datos: 1 PRIMARY + 2 SECONDARY."),
    ("Cómo funciona:", "si el PRIMARY cae, eligen uno nuevo solos; cada pedido + descuento de stock es una transacción atómica."),
    ("Por qué:", "los pedidos son dinero — no se pueden perder ni quedar a medias."),
], "DEMO: voy a matar un nodo y el sistema sigue vendiendo.")

content_slide("NetworkPolicy — firewall entre pods", [
    ("Qué es:", "reglas de quién puede hablar con quién dentro del cluster."),
    ("Cómo funciona:", "default-deny (bloqueo todo) + permitir solo lo necesario. Lo aplica el CNI Calico."),
    ("Por qué:", "si un pod se compromete, no puede moverse hasta la base de datos."),
], "DEMO: un pod sin permiso intenta conectarse a Mongo y la conexión se cae: bloqueado por política.")

content_slide("KEDA — escalado automático por demanda", [
    ("Qué es:", "autoscaler que mide la longitud de la cola Redis del worker."),
    ("Cómo funciona:", "si se acumulan pedidos, KEDA crea más pods del worker (1→5); cuando la cola se vacía, los baja a 1."),
    ("Por qué:", "aguanta picos sin desperdiciar recursos en horas muertas."),
], "DEMO: inyecto carga a la cola y verán al worker pasar de 1 a 4 pods solo.")

content_slide("CI/CD — Tekton + ArgoCD (GitOps)", [
    ("Tekton:", "al hacer git push, construye la imagen (Kaniko) y la sube al registro local Y a Docker Hub."),
    ("ArgoCD:", "vigila Git y despliega solo; con selfHeal revierte cambios manuales → Git es la única verdad."),
    ("Por qué:", "despliegues reproducibles, auditables y sin tocar el cluster a mano."),
], "Hago un push y el cluster se actualiza solo; si alguien edita a mano, ArgoCD lo regresa al estado de Git.")

content_slide("Argo Rollouts — Canary (cero downtime)", [
    ("Qué es:", "despliegue gradual de versiones nuevas del API."),
    ("Cómo funciona:", "la versión nueva recibe 10% → 50% → 100% por pasos; un AnalysisTemplate hace rollback automático si el error sube."),
    ("Por qué:", "una versión rota no tumba a todos; se detecta y revierte sola."),
], "DEMO: promuevo el canary en vivo; si el error-rate sube, el rollout se devuelve solo.")

content_slide("Observabilidad — Loki + Prometheus + Grafana", [
    ("Métricas (Prometheus):", "latencia del API, longitud de cola, error-rate + alertas automáticas."),
    ("Logs (Loki + Alloy):", "logs centralizados; busco un pedido por su pedidoId en segundos."),
    ("Por qué:", "en producción, sin medir no operas — detectas problemas antes que el cliente."),
], "DEMO: busco un pedido por su ID en los logs de Loki en vivo.")

content_slide("Seguridad y gobernanza", [
    ("Ingress + TLS:", "acceso HTTPS (zuyu.local) con certificado de cert-manager (CA propia)."),
    ("ResourceQuota:", "límites de CPU/memoria por namespace → ningún servicio acapara el cluster."),
    ("Secrets:", "cifrados con SealedSecrets; contraseñas con Argon2."),
], "DEMO: intento crear un pod que excede la cuota y el cluster lo rechaza.")

table_slide("Stack en una mirada",
    ["Capa", "Herramienta"],
    [
        ["Aplicación", "Node.js + Express + EJS + Alpine.js"],
        ["Datos", "MongoDB 7 (Replica Set) · Redis 7 + BullMQ"],
        ["Cluster", "kubeadm + Calico (VXLAN), 3 nodos"],
        ["CI/CD", "Tekton + ArgoCD + Argo Rollouts"],
        ["Escalado", "KEDA + HPA"],
        ["Observabilidad", "Prometheus + Loki + Grafana"],
        ["Seguridad", "NetworkPolicy + cert-manager + SealedSecrets"],
    ],
    "Todo esto es estándar de la industria; lo monté de cero sobre bare-metal.")

closing_slide("Gracias", [
    "github.com/SEBASTIANCONTRERAS35/micrositio-pedidos",
    "¿Preguntas?  →  pasamos a las 7 demos en el cluster",
], "Pasar directo a las demos en vivo.")

OUT = "docs/presentacion-ejecutiva.pptx"
prs.save(OUT)
print(f"OK: {OUT} — {len(prs.slides._sldIdLst)} slides")
